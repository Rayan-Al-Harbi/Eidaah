import os
import torch
import gradio as gr
from typing import List
from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document
from pptx import Presentation
from difflib import SequenceMatcher
from transformers import pipeline, AutoTokenizer, AutoModelForCausalLM
import warnings

warnings.filterwarnings("ignore")

# Load fast Qwen model on CPU or GPU
def get_llm():
    model_name = "Qwen/Qwen2-1.5B-Instruct"
    print(f"Loading optimized model: {model_name}")

    device = "cuda" if torch.cuda.is_available() else "cpu"
    print(f"Using device: {device}")

    tokenizer = AutoTokenizer.from_pretrained(model_name)

    # ✅ FIX: Conditional device_map based on available device
    if device == "cuda":
        model = AutoModelForCausalLM.from_pretrained(
            model_name,
            device_map={"": 0}  # GPU: keep whole model on GPU
        )
    else:
        # CPU: Don't use device_map, use standard loading
        model = AutoModelForCausalLM.from_pretrained(
            model_name,
            torch_dtype=torch.float32  # Use full precision on CPU
        )
        model = model.to(device)  # Move to CPU explicitly

    pipe = pipeline(
        "text-generation",
        model=model,
        tokenizer=tokenizer,
        max_new_tokens=60,  
        do_sample=False,
        return_full_text=False,
        device=0 if device == "cuda" else -1  # 0 for GPU, -1 for CPU
    )

    return pipe, tokenizer

# Initialize model globally
print("Initializing AI model... This may take a few minutes on first run.")
llm_pipe, llm_tokenizer = get_llm()
print("✅ Model loaded successfully!")


# Universal Loader for PDF/PPTX
def load_document(file_path: str) -> List[Document]:
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".pdf":
            loader = PyPDFLoader(file_path)
            return loader.load()
        elif ext == ".pptx":
            prs = Presentation(file_path)
            slides = []
            for i, slide in enumerate(prs.slides):
                text = "\n".join(
                    shape.text for shape in slide.shapes if hasattr(shape, "text")
                )
                slides.append(Document(page_content=text.strip(), metadata={"page": i + 1}))
            return slides
        else:
            return []
    except Exception as e:
        print(f"Error loading document: {e}")
        return []

# (MAIN METHOD) Two-step generation
def generate_explanation_and_example(text: str):
    if not text.strip():
        return "No text found on this slide.", "No example available."
    try:
        # Step 1 – Explanation
        explain_prompt = f"""Analyze and explain this presentation slide content clearly and concisely:

{text}

Provide a brief 2-3 sentence explanation of what this slide is about."""
        messages = [{"role": "user", "content": explain_prompt}]
        formatted_prompt = llm_tokenizer.apply_chat_template(messages, tokenize=False, add_generation_prompt=True)
        explanation = llm_pipe(formatted_prompt, max_new_tokens=80)[0]["generated_text"].strip().lstrip(': \n')

        # Step 2 – Example
        example_prompt = f"""Based on this topic: {text[:200]}

Give ONE concrete, practical example that illustrates this concept. Keep it brief and specific."""
        example_messages = [{"role": "user", "content": example_prompt}]
        formatted_example_prompt = llm_tokenizer.apply_chat_template(example_messages, tokenize=False, add_generation_prompt=True)
        example = llm_pipe(formatted_example_prompt, max_new_tokens=60)[0]["generated_text"].strip().lstrip(': \n')

        # Prevent copy-paste repetition
        similarity = SequenceMatcher(None, explanation.lower(), example.lower()).ratio()
        if similarity > 0.8:
            example = "Example: (could not generate a unique one). Try a different slide."

        return explanation, example
        
    except Exception as e:
        print(f"Error in generate_explanation_and_example: {e}")
        return f"Error: {e}", ""